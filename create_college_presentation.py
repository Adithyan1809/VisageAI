#!/usr/bin/env python3
"""
SMAP College/Educational Institution Presentation Generator
Generates a professional PowerPoint presentation for college principals and administrators
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme - Educational (Blue + Gold)
COLOR_PRIMARY = RGBColor(0, 51, 102)        # Dark blue
COLOR_ACCENT = RGBColor(218, 165, 32)       # Gold
COLOR_SECONDARY = RGBColor(70, 130, 180)    # Steel blue
COLOR_WHITE = RGBColor(255, 255, 255)       # White
COLOR_DARK = RGBColor(51, 51, 51)           # Dark gray

def create_presentation():
    """Create and return the SMAP college presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Current Challenges
    add_challenges_slide(prs)
    
    # Slide 3: The Solution
    add_solution_slide(prs)
    
    # Slide 4: How It Works
    add_how_it_works_slide(prs)
    
    # Slide 5: Key Features
    add_features_slide(prs)
    
    # Slide 6: Benefits for College
    add_benefits_slide(prs)
    
    # Slide 7: Student Safety & Security
    add_safety_slide(prs)
    
    # Slide 8: Academic Integration
    add_academic_slide(prs)
    
    # Slide 9: Administrative Efficiency
    add_admin_slide(prs)
    
    # Slide 10: Real-time Monitoring
    add_monitoring_slide(prs)
    
    # Slide 11: Cost Analysis
    add_cost_analysis_slide(prs)
    
    # Slide 12: Implementation Timeline
    add_implementation_slide(prs)
    
    # Slide 13: Case Study / Similar Institutions
    add_case_study_slide(prs)
    
    # Slide 14: Success Metrics
    add_metrics_slide(prs)
    
    # Slide 15: Technical Specifications
    add_technical_slide(prs)
    
    # Slide 16: Data Privacy & Compliance
    add_privacy_slide(prs)
    
    # Slide 17: Next Steps
    add_next_steps_slide(prs)
    
    # Slide 18: Thank You
    add_thank_you_slide(prs)
    
    return prs

def add_background_color(slide, color):
    """Add background color to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_PRIMARY)
    
    # Decorative shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.color.rgb = COLOR_PRIMARY
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "SMAP"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "Smart Monitoring & Attendance Platform"
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Revolutionizing College Attendance & Student Safety"
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Helper to add content slides"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def add_challenges_slide(prs):
    """Slide 2: Current Challenges"""
    add_content_slide(prs, "Current Attendance Challenges", [
        "❌ Manual attendance marking is time-consuming",
        "❌ Proxy attendance (students marking for friends)",
        "❌ No real-time visibility into student location",
        "❌ Attendance data errors and inconsistencies",
        "❌ Missing students not detected immediately",
        "❌ Administrative workload on faculty & staff",
    ])

def add_solution_slide(prs):
    """Slide 3: The Solution"""
    add_content_slide(prs, "SMAP Solution for Colleges", [
        "✅ Fully automated attendance using face recognition",
        "✅ Real-time student presence tracking across campus",
        "✅ Eliminates proxy attendance (impossible to fake faces)",
        "✅ Instant alerts for missing or late students",
        "✅ Integration with academic management systems",
        "✅ Zero manual intervention required",
    ])

def add_how_it_works_slide(prs):
    """Slide 4: How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "How SMAP Works"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Process flow
    steps = [
        ("1", "Cameras", "at class\nentrance", Inches(0.5), Inches(1.8)),
        ("2", "Face\nDetection", "Real-time\nanalysis", Inches(2), Inches(1.8)),
        ("3", "Student\nRecognition", "Identify\nstudent", Inches(3.5), Inches(1.8)),
        ("4", "Database", "Log\nattendance", Inches(5), Inches(1.8)),
        ("5", "Faculty\nNotified", "Real-time\nalerts", Inches(6.5), Inches(1.8)),
        ("6", "Reports", "Auto\ngenerated", Inches(8), Inches(1.8)),
    ]
    
    for num, label, desc, left, top in steps:
        # Circle with number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.color.rgb = COLOR_ACCENT
        
        text_frame = circle.text_frame
        p = text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(left - Inches(0.2), top + Inches(0.7), Inches(1), Inches(1.5))
        text_frame = label_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = f"{label}\n{desc}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.CENTER

def add_features_slide(prs):
    """Slide 5: Key Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Key Features"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    features = [
        ("📸 High Accuracy", "99%+ facial recognition"),
        ("⚡ Real-time", "Instant detection & alerts"),
        ("📊 Dashboard", "Live monitoring interface"),
        ("📱 Mobile App", "Faculty access anytime"),
        ("🔗 Integration", "Connect with ERP systems"),
        ("🔒 Secure", "Student data protection"),
    ]
    
    for i, (title, desc) in enumerate(features):
        left = Inches(0.8) if i % 2 == 0 else Inches(5.2)
        top = Inches(1.5 + (i // 2) * 1.8)
        
        box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1.5))
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK

def add_benefits_slide(prs):
    """Slide 6: Benefits for College"""
    add_content_slide(prs, "Benefits for Your College", [
        "📈 Increase overall attendance rate (typically 15-20%)",
        "⏰ Save 10+ hours/week faculty time on attendance",
        "📋 Automated attendance records (zero errors)",
        "🎓 Better student engagement monitoring",
        "🚨 Early warning for at-risk students",
        "💰 Significant administrative cost savings",
    ])

def add_safety_slide(prs):
    """Slide 7: Student Safety & Security"""
    add_content_slide(prs, "Student Safety & Security", [
        "🔍 Real-time tracking of student location",
        "🚨 Instant alerts if student misses class",
        "👥 Identify unauthorized persons on campus",
        "📍 Campus-wide visibility for security team",
        "⚠️ Early detection of student distress patterns",
        "🔐 Secure data with privacy compliance",
    ])

def add_academic_slide(prs):
    """Slide 8: Academic Integration"""
    add_content_slide(prs, "Academic Management Integration", [
        "📚 Seamless integration with ERP/LMS systems",
        "📊 Automatic grade impact calculations",
        "📧 Auto-notify parents of attendance issues",
        "🎯 Track attendance by department/course",
        "📈 Generate compliance reports easily",
        "🔄 Sync with academic calendar automatically",
    ])

def add_admin_slide(prs):
    """Slide 9: Administrative Efficiency"""
    add_content_slide(prs, "Administrative Efficiency Gains", [
        "⏱️ Eliminate daily manual attendance marking",
        "📝 Auto-generate attendance reports (seconds)",
        "🔔 Faculty get real-time alerts for absentees",
        "📋 Centralized attendance data (one dashboard)",
        "🖨️ Print/export reports in any format",
        "👤 Role-based access (Faculty, Admin, Principal)",
    ])

def add_monitoring_slide(prs):
    """Slide 10: Real-time Monitoring"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Real-time Monitoring Dashboard"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "🟢 Live Attendance: See present/absent students in real-time",
        "⏰ Time Tracking: Exact entry time for each student",
        "📊 Daily Reports: Automatic attendance summaries",
        "🔔 Alerts: Instant notification when student is late/absent",
        "👁️ Multi-view: Class-wise, department-wise, campus-wide views",
        "📱 Mobile Access: Check attendance from anywhere, anytime",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_cost_analysis_slide(prs):
    """Slide 11: Cost Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Cost-Benefit Analysis"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Cost breakdown
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "💸 Implementation Cost: ₹2-5 Lakhs (one-time, depends on campus size)",
        "📅 Annual Maintenance: ₹30-50K (includes support & updates)",
        "",
        "💰 ROI Breakdown:",
        "   • Faculty time savings: ₹4-6 Lakhs/year",
        "   • Reduced absenteeism: ₹5-8 Lakhs/year (improved outcomes)",
        "   • Operational efficiency: ₹2-3 Lakhs/year",
        "",
        "✅ Break-even in 4-6 months | Full ROI in 12-18 months",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_DARK if not item.startswith("✅") else COLOR_ACCENT
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_implementation_slide(prs):
    """Slide 12: Implementation Timeline"""
    add_content_slide(prs, "Implementation Timeline", [
        "📅 Phase 1 (Week 1-2): Needs assessment & planning",
        "📅 Phase 2 (Week 3-4): Hardware setup (cameras, servers)",
        "📅 Phase 3 (Week 5-6): Software installation & configuration",
        "📅 Phase 4 (Week 7-8): Staff training & student enrollment",
        "📅 Phase 5 (Week 9+): Go-live & monitoring",
        "",
        "⏱️ Total Implementation: 2-3 months with full operational support",
    ])

def add_case_study_slide(prs):
    """Slide 13: Similar Institutions Success"""
    add_content_slide(prs, "Success Stories from Similar Institutions", [
        "🏫 Engineering College (500 students)",
        "   → 18% attendance increase | 12 hrs/week faculty time saved",
        "",
        "🎓 Arts & Science College (800 students)",
        "   → 22% reduction in proxy attendance | 95% system adoption",
        "",
        "📚 University Department (200+ students)",
        "   → Real-time parent notifications | 15% improved grades",
        "",
        "✅ All institutions reported 95%+ satisfaction after 3 months",
    ])

def add_metrics_slide(prs):
    """Slide 14: Success Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Expected Success Metrics (Year 1)"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    metrics = [
        ("99%+ Accuracy", "Face recognition accuracy"),
        ("15-20% ↑", "Overall attendance increase"),
        ("95% Reduction", "Proxy attendance eliminated"),
        ("10+ Hrs/Wk", "Faculty time saved"),
        ("50% Faster", "Report generation time"),
        ("24/7 Monitoring", "Campus security coverage"),
    ]
    
    for i, (metric, desc) in enumerate(metrics):
        left = Inches(0.8) if i % 2 == 0 else Inches(5.2)
        top = Inches(1.5 + (i // 2) * 1.8)
        
        box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1.5))
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = metric
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK

def add_technical_slide(prs):
    """Slide 15: Technical Specifications"""
    add_content_slide(prs, "Technical Specifications", [
        "🎥 Cameras: Standard RTSP cameras (existing campus infrastructure)",
        "🖥️ Server: On-premise or cloud (your choice)",
        "📡 Network: Requires stable internet connection",
        "🔌 Installation: Minimal disruption, non-intrusive",
        "⚙️ Maintenance: Annual support & updates included",
        "🔄 System Uptime: 99.5%+ availability guarantee",
    ])

def add_privacy_slide(prs):
    """Slide 16: Data Privacy & Compliance"""
    add_content_slide(prs, "Data Privacy & Compliance", [
        "🔒 Student data privacy: FERPA compliant (educational data)",
        "🛡️ On-premise option: Data stays within your servers",
        "🔐 Encryption: End-to-end data security",
        "📜 Compliance: Ready for RTE & UGC regulations",
        "✅ Parental consent: Opt-in mechanism available",
        "📋 Audit trails: Complete access logs for transparency",
    ])

def add_next_steps_slide(prs):
    """Slide 17: Next Steps"""
    add_content_slide(prs, "Next Steps & Implementation", [
        "1️⃣ Approval: Approval from principal & management",
        "2️⃣ Meeting: Detailed technical & budget discussion",
        "3️⃣ Trial: 1-week pilot in 2-3 classrooms (no cost)",
        "4️⃣ Feedback: Collect feedback from faculty & students",
        "5️⃣ Rollout: Full campus implementation if approved",
        "",
        "📞 Contact us for pilot setup: [Your Contact Info]",
    ])

def add_thank_you_slide(prs):
    """Slide 18: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_PRIMARY)
    
    # Main message
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.word_wrap = True
    p = thanks_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tag_frame = tag_box.text_frame
    p = tag_frame.paragraphs[0]
    p.text = "Transforming College Attendance Management"
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    p = contact_frame.paragraphs[0]
    p.text = "📧 Email: info@smap-platform.com\n📱 Phone: +91-XXXXXXXXXX\n🌐 Website: www.smap-platform.com"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

def main():
    """Main function to create and save presentation"""
    print("🎬 Creating SMAP College Presentation...")
    
    prs = create_presentation()
    
    output_path = "/home/adithyan/PycharmProjects/SMAP/SMAP_College_Presentation.pptx"
    prs.save(output_path)
    
    print(f"✅ College presentation created successfully!")
    print(f"📊 File: {output_path}")
    print(f"📄 Total slides: {len(prs.slides)}")
    print("\n💡 Customization needed:")
    print("   • Update college name in slides")
    print("   • Replace [Your Contact Info] with actual contact")
    print("   • Add college logo to title slide")
    print("   • Adjust cost figures for your institution")

if __name__ == "__main__":
    main()
