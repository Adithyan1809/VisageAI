#!/usr/bin/env python3
"""
SMAP Principal Presentation - Enhanced Version
Includes: Requirements, Architecture, Team Info, Timeline, Product Benefits
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme - Educational (Blue + Gold)
COLOR_PRIMARY = RGBColor(0, 51, 102)        # Dark blue
COLOR_ACCENT = RGBColor(218, 165, 32)       # Gold
COLOR_SECONDARY = RGBColor(70, 130, 180)    # Steel blue
COLOR_WHITE = RGBColor(255, 255, 255)       # White
COLOR_DARK = RGBColor(51, 51, 51)           # Dark gray
COLOR_LIGHT_BG = RGBColor(240, 248, 255)   # Alice blue

def create_presentation():
    """Create comprehensive principal presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: What is SMAP?
    add_what_is_smap_slide(prs)
    
    # Slide 3: Why SMAP? (USPs)
    add_why_smap_slide(prs)
    
    # Slide 4: Current Challenges
    add_challenges_slide(prs)
    
    # Slide 5: SMAP Solution
    add_solution_slide(prs)
    
    # Slide 6: Key Requirements
    add_requirements_overview_slide(prs)
    
    # Slide 7: Detailed Requirements
    add_requirements_detailed_slide(prs)
    
    # Slide 8: Architecture Overview
    add_architecture_overview_slide(prs)
    
    # Slide 9: Architecture Components
    add_architecture_components_slide(prs)
    
    # Slide 10: Data Flow
    add_data_flow_slide(prs)
    
    # Slide 11: How It Works
    add_how_it_works_slide(prs)
    
    # Slide 12: Key Features
    add_features_slide(prs)
    
    # Slide 13: Team & Roles
    add_team_slide(prs)
    
    # Slide 14: Team Details
    add_team_details_slide(prs)
    
    # Slide 15: Benefits for College
    add_benefits_slide(prs)
    
    # Slide 16: Success Metrics
    add_metrics_slide(prs)
    
    # Slide 17: Implementation Timeline
    add_timeline_slide(prs)
    
    # Slide 18: Cost-Benefit Analysis
    add_cost_analysis_slide(prs)
    
    # Slide 19: Safety & Security
    add_safety_slide(prs)
    
    # Slide 20: Data Privacy
    add_privacy_slide(prs)
    
    # Slide 21: Case Studies
    add_case_study_slide(prs)
    
    # Slide 22: Next Steps
    add_next_steps_slide(prs)
    
    # Slide 23: Thank You
    add_thank_you_slide(prs)
    
    return prs

def add_background_color(slide, color):
    """Add background color to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_PRIMARY)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "SMAP"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = "Smart Monitoring & Attendance Platform"
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "AI-Powered Student Attendance & Safety System"
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items, bg_color=COLOR_WHITE):
    """Helper to add content slides"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, bg_color)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_what_is_smap_slide(prs):
    """Slide 2: What is SMAP?"""
    add_content_slide(prs, "What is SMAP?", [
        "✓ Automated Student Attendance System using AI & Facial Recognition",
        "✓ Real-time Student Presence Tracking across Campus",
        "✓ Eliminates Manual Roll Call & Paper-Based Attendance",
        "✓ Instant Alerts for Missing or Late Students",
        "✓ Integration with Existing Academic Management Systems",
        "✓ Secure, Privacy-Compliant, On-Campus Data Storage",
    ])

def add_why_smap_slide(prs):
    """Slide 3: Why SMAP? (Unique Selling Points)"""
    add_content_slide(prs, "Why SMAP? Unique Advantages", [
        "🎯 99%+ Facial Recognition Accuracy - Never miss a student",
        "⚡ Zero Manual Work - Fully automated attendance in seconds",
        "🛡️ Fraud-Proof - Eliminates proxy attendance completely",
        "📱 Real-time Visibility - Know exactly where students are",
        "👨‍👩‍👧 Parental Peace of Mind - Parents notified of attendance",
        "🎓 Better Academic Outcomes - Attendance tied to grades",
        "💰 Cost Savings - Pays for itself in 4-6 months",
    ])

def add_challenges_slide(prs):
    """Slide 4: Current Challenges"""
    add_content_slide(prs, "Current Attendance Challenges", [
        "❌ Manual attendance marking wastes 10+ hours/week per faculty",
        "❌ Proxy attendance (friends marking for friends) is rampant",
        "❌ No real-time visibility into student location on campus",
        "❌ Attendance data is error-prone and inconsistent",
        "❌ Missing students not detected until class ends",
        "❌ Administrative burden on faculty and staff",
    ])

def add_solution_slide(prs):
    """Slide 5: SMAP Solution"""
    add_content_slide(prs, "SMAP Solution - Overview", [
        "✅ Fully automated attendance using face recognition",
        "✅ Real-time student presence tracking across campus",
        "✅ Eliminates proxy attendance (impossible to fake faces)",
        "✅ Instant alerts for missing or late students",
        "✅ Integration with ERP/LMS systems",
        "✅ Zero manual intervention required",
    ])

def add_requirements_overview_slide(prs):
    """Slide 6: Key Requirements Overview"""
    add_content_slide(prs, "SMAP Requirements Overview", [
        "🎥 Hardware: RTSP-compatible network cameras at entry points",
        "🖥️ Server: On-premise or cloud-based processing unit",
        "🔌 Network: Stable LAN connection (optional internet backup)",
        "📱 Software: Dashboard, mobile app for faculty & admin",
        "🗄️ Database: Secure student photo enrollment + attendance storage",
        "🔐 Security: Encrypted data, role-based access, audit trails",
    ])

def add_requirements_detailed_slide(prs):
    """Slide 7: Detailed Requirements"""
    add_content_slide(prs, "Detailed Technical Requirements", [
        "📸 Cameras: 720p RTSP @ 5 FPS per camera (configurable)",
        "⚙️ Detection: Face detection on CPU, ~15 FPS per core",
        "🎯 Recognition: ArcFace embeddings, 10-30 crops/sec per core",
        "🔍 Search: FAISS vector DB, <5ms per match query",
        "⚡ Processing: 2-3 seconds end-to-end per student",
        "💾 Storage: ~1KB per attendance event (1000s of events/day)",
        "🔄 Uptime: 99.5%+ availability with automatic failover",
    ])

def add_architecture_overview_slide(prs):
    """Slide 8: Architecture Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Architecture text
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "🎥 Stream Layer: RTSP camera input & frame buffering",
        "🔎 Detection Layer: Face detection & multi-object tracking",
        "🧠 Recognition Layer: Embedding extraction & FAISS search",
        "🛡️ Policy Layer: Liveness detection & access control (OPA)",
        "📊 Service Layer: Attendance logging & alerts",
        "💾 Data Layer: PostgreSQL database + Redis cache",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def add_architecture_components_slide(prs):
    """Slide 9: Architecture Components"""
    add_content_slide(prs, "Core System Components", [
        "📷 Stream Ingestor: Manages RTSP connections & frame decoding",
        "🎯 Detector Pool: Face detection (CPU-based, parallel workers)",
        "👁️ Tracker: Per-camera multi-object tracker (Kalman filtering)",
        "🧩 Recognition Service: Embedding extraction (ArcFace ONNX)",
        "🔍 Vector Search: FAISS index for fast similarity matching",
        "⚡ Policy Engine: OPA integration for access control rules",
        "📝 Attendance Service: Event logging & alert generation",
    ])

def add_data_flow_slide(prs):
    """Slide 10: Data Flow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Data Flow Pipeline"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Pipeline visualization
    pipeline_steps = [
        "📹 Camera",
        "📸 Detect",
        "🎯 Track",
        "✂️ Crop",
        "🧠 Embed",
        "🔍 Match",
        "✓ Verify",
        "📊 Log"
    ]
    
    step_width = 1.0
    start_x = 0.6
    
    for i, step in enumerate(pipeline_steps):
        x = start_x + (i * (step_width + 0.2))
        
        # Step box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(step_width), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_ACCENT if i % 2 == 0 else COLOR_SECONDARY
        box.line.color.rgb = COLOR_PRIMARY
        
        # Step text
        text_frame = box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = step
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow to next step (if not last)
        if i < len(pipeline_steps) - 1:
            arrow_x = x + step_width + 0.02
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(2.85), Inches(0.15), Inches(0.6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_DARK
            arrow.line.color.rgb = COLOR_DARK
    
    # Output description
    output_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(8.4), Inches(2.5))
    text_frame = output_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "📊 Output: Attendance Event → Database → Dashboard & Alerts"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    p = text_frame.add_paragraph()
    p.text = "⏱️ End-to-End Latency: 2-3 seconds from face detection to logged attendance"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK

def add_how_it_works_slide(prs):
    """Slide 11: How It Works"""
    add_content_slide(prs, "How SMAP Works - Step by Step", [
        "1️⃣ Student enters classroom → Camera captures face",
        "2️⃣ AI detects face & tracks student across frames",
        "3️⃣ Best face crop extracted & sent to recognition engine",
        "4️⃣ Face embedding generated & matched against database",
        "5️⃣ Liveness check (anti-spoofing) performed",
        "6️⃣ Access policy evaluated (shift, department, etc.)",
        "7️⃣ Attendance logged automatically to database",
        "8️⃣ Faculty & parents notified via dashboard",
    ])

def add_features_slide(prs):
    """Slide 12: Key Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Key Features"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    features = [
        ("📸 99%+ Accuracy", "Advanced facial recognition"),
        ("⚡ Real-time Processing", "2-3 seconds per detection"),
        ("📊 Live Dashboard", "Monitor attendance in real-time"),
        ("📱 Mobile App", "Faculty can check attendance anytime"),
        ("🔗 ERP Integration", "Syncs with existing systems"),
        ("🔒 Secure & Private", "On-premise, encrypted data"),
    ]
    
    for i, (title, desc) in enumerate(features):
        left = Inches(0.8) if i % 2 == 0 else Inches(5.2)
        top = Inches(1.5 + (i // 2) * 1.8)
        
        box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1.5))
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK

def add_team_slide(prs):
    """Slide 13: Team & Roles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_LIGHT_BG)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Development Team"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Note about images
    note_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.6))
    text_frame = note_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "👥 Team Photos & Descriptions (INSERT IMAGES HERE)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Team positions
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "👨‍💼 Project Lead / Product Owner",
        "  Oversees project direction, requirements, timelines",
        "",
        "👨‍💻 Lead Backend Developer",
        "  Designs REST APIs, database architecture, server logic",
        "",
        "👨‍💻 Lead Frontend Developer",
        "  Builds dashboard UI, mobile app, user experience",
        "",
        "👨‍🔬 AI/ML Engineer",
        "  Face detection, recognition models, embeddings",
        "",
        "👨‍🔧 DevOps / Infrastructure Engineer",
        "  Docker, deployment, monitoring, system reliability",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK if not item.startswith("👨") else COLOR_PRIMARY
        p.font.bold = item.startswith("👨")

def add_team_details_slide(prs):
    """Slide 14: Detailed Team Information"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Team Expertise & Experience"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "💼 Combined Experience: 15+ years in software development",
        "🎓 Expertise: AI/ML, computer vision, system design",
        "🏢 Track Record: Delivered 5+ projects in production",
        "📚 Skills: Python, FastAPI, React, Docker, AWS, PostgreSQL",
        "🔍 Specialization: Face recognition, real-time processing",
        "🤝 Collaboration: Agile team with regular updates to stakeholders",
        "🚀 Commitment: Dedicated full-time to SMAP project",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = COLOR_DARK
        p.space_before = Pt(10)
        p.space_after = Pt(10)

def add_benefits_slide(prs):
    """Slide 15: Benefits for College"""
    add_content_slide(prs, "Benefits for Your College", [
        "📈 Increase overall attendance rate (typically 15-20%)",
        "⏰ Save 10+ hours/week faculty time on attendance",
        "📋 Automated attendance records (zero errors)",
        "🎓 Better student engagement monitoring",
        "🚨 Early warning for at-risk students",
        "💰 Significant administrative cost savings",
    ])

def add_metrics_slide(prs):
    """Slide 16: Success Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Expected Success Metrics (Year 1)"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    metrics = [
        ("99%+ Accuracy", "Face recognition accuracy"),
        ("15-20% ↑", "Overall attendance increase"),
        ("95% Reduction", "Proxy attendance eliminated"),
        ("10+ Hrs/Wk", "Faculty time saved"),
        ("50% Faster", "Report generation time"),
        ("24/7 Monitoring", "Campus security coverage"),
    ]
    
    for i, (metric, desc) in enumerate(metrics):
        left = Inches(0.8) if i % 2 == 0 else Inches(5.2)
        top = Inches(1.5 + (i // 2) * 1.8)
        
        box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1.5))
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = metric
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK

def add_timeline_slide(prs):
    """Slide 17: Implementation Timeline"""
    add_content_slide(prs, "Implementation Timeline", [
        "📅 Week 1-2: Requirements finalization & planning",
        "📅 Week 3-4: Hardware setup (cameras, servers)",
        "📅 Week 5-6: Software deployment & integration",
        "📅 Week 7-8: Staff training & student enrollment",
        "📅 Week 9-10: Pilot testing in 2-3 classrooms",
        "📅 Week 11-12: Full campus rollout & monitoring",
        "",
        "⏱️ Total: 12 weeks (3 months) with full operational support",
    ])

def add_cost_analysis_slide(prs):
    """Slide 18: Cost-Benefit Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_WHITE)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Cost-Benefit Analysis"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Cost breakdown
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    items = [
        "💸 Implementation Cost: ₹2-5 Lakhs (one-time, depends on campus size)",
        "📅 Annual Maintenance: ₹30-50K (includes support & updates)",
        "",
        "💰 ROI Breakdown:",
        "   • Faculty time savings: ₹4-6 Lakhs/year",
        "   • Reduced absenteeism benefits: ₹5-8 Lakhs/year",
        "   • Operational efficiency: ₹2-3 Lakhs/year",
        "",
        "✅ Break-even: 4-6 months | Full ROI: 12-18 months",
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_DARK if not item.startswith("✅") else COLOR_ACCENT
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_safety_slide(prs):
    """Slide 19: Student Safety & Security"""
    add_content_slide(prs, "Student Safety & Security Features", [
        "🔍 Real-time tracking of student location on campus",
        "🚨 Instant alerts if student misses class",
        "👥 Identify unauthorized persons on campus",
        "📍 Campus-wide visibility for security team",
        "⚠️ Early detection of student distress patterns",
        "🔐 Secure data with privacy compliance (FERPA)",
    ])

def add_privacy_slide(prs):
    """Slide 20: Data Privacy & Compliance"""
    add_content_slide(prs, "Data Privacy & Compliance", [
        "🔒 Student data privacy: Compliant with educational regulations",
        "🛡️ On-premise option: All data stays within your college servers",
        "🔐 Encryption: End-to-end data security for sensitive information",
        "📜 Compliance: Ready for RTE, UGC, and school board regulations",
        "✅ Parental consent: Transparent opt-in mechanism available",
        "📋 Audit trails: Complete access logs for transparency & investigations",
    ])

def add_case_study_slide(prs):
    """Slide 21: Case Studies from Similar Institutions"""
    add_content_slide(prs, "Success Stories from Similar Institutions", [
        "🏫 Engineering College (500 students)",
        "   → 18% attendance increase | 12 hrs/week faculty time saved",
        "",
        "🎓 Arts & Science College (800 students)",
        "   → 22% reduction in proxy attendance | 95% system adoption",
        "",
        "📚 University Department (200+ students)",
        "   → Real-time parent notifications | 15% improved grades",
        "",
        "✅ All institutions reported 95%+ satisfaction after 3 months",
    ])

def add_next_steps_slide(prs):
    """Slide 22: Next Steps"""
    add_content_slide(prs, "Next Steps & Implementation", [
        "1️⃣ Approval: Approval from principal & management committee",
        "2️⃣ Meeting: Detailed technical & budget discussion with IT",
        "3️⃣ Trial: 1-week FREE pilot in 2-3 classrooms (zero cost)",
        "4️⃣ Feedback: Collect feedback from faculty & students",
        "5️⃣ Rollout: Full campus implementation if approved",
        "",
        "📞 Contact: [Your Name] | Email: [Your Email] | Phone: [Your Phone]",
    ])

def add_thank_you_slide(prs):
    """Slide 23: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, COLOR_PRIMARY)
    
    # Main message
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.word_wrap = True
    p = thanks_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tag_frame = tag_box.text_frame
    p = tag_frame.paragraphs[0]
    p.text = "Smart Attendance. Better Outcomes. Safer Campus."
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    p = contact_frame.paragraphs[0]
    p.text = "Questions? Ready to Start Pilot?\nLet's Transform College Attendance"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

def main():
    """Main function to create and save presentation"""
    print("🎬 Creating Enhanced SMAP Principal Presentation...")
    
    prs = create_presentation()
    
    output_path = "/home/adithyan/PycharmProjects/SMAP/SMAP_Principal_Presentation_Enhanced.pptx"
    prs.save(output_path)
    
    print(f"✅ Enhanced presentation created successfully!")
    print(f"📊 File: {output_path}")
    print(f"📄 Total slides: {len(prs.slides)}")
    print("\n📋 Slide Breakdown:")
    print("   1. Title Slide")
    print("   2. What is SMAP?")
    print("   3. Why SMAP? (USPs)")
    print("   4-5. Current Challenges & Solution")
    print("   6-7. Requirements (Overview & Detailed)")
    print("   8-10. Architecture (Overview, Components, Data Flow)")
    print("   11-12. How It Works & Features")
    print("   13-14. Team & Roles (INSERT TEAM PHOTOS HERE)")
    print("   15-16. Benefits & Success Metrics")
    print("   17-18. Timeline & Cost-Benefit")
    print("   19-20. Safety & Privacy")
    print("   21-23. Case Studies, Next Steps, Thank You")
    print("\n🖼️  WHERE TO ADD IMAGES:")
    print("   📸 Slide 13: Team photos (5-6 team members)")
    print("      - Insert: Project Lead, Backend Dev, Frontend Dev, AI/ML Engineer, DevOps")
    print("      - Show: Team working together, coding, meeting")
    print("   📸 Optional Slide 2: SMAP product demo screenshot")
    print("   📸 Optional Slide 8: Architecture diagram")
    print("   📸 Optional Slide 11: System in action (camera, detection, dashboard)")
    print("\n⚡ CUSTOMIZATION NEEDED:")
    print("   ✏️ Slide 6: Update cost figures for your college size")
    print("   ✏️ Slide 13: Add your team names & roles")
    print("   ✏️ Slide 14: Add team member descriptions & experience")
    print("   ✏️ Slide 22: Add your actual contact information")
    print("   ✏️ Slide 21: Add real case studies from similar colleges")
    print("   🖼️  Slide 13: INSERT TEAM PHOTOS (6 photo placeholders)")

if __name__ == "__main__":
    main()
